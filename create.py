from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE

#new presentation
prs = Presentation()

#title slide
slide_layout = prs.slide_layouts[0]
slide1 = prs.slides.add_slide(slide_layout)
subtitle = slide1.placeholders[1]
title = slide1.shapes.title
fill1 = slide1.background.fill
fill1.solid()
fill1.fore_color.rgb = RGBColor(173,216,230) #light blue
title.text = "Powerpoint from Python"
subtitle.text = "Built for Powerpynt with python-pptx"

#company slide
slide_layout = prs.slide_layouts[0]
slide2 = prs.slides.add_slide(slide_layout)
subtitle2 = slide2.placeholders[1]
title2 = slide2.shapes.title
fill2 = slide2.background.fill
fill2.solid()
fill2.fore_color.rgb = RGBColor(224, 255, 255)
title2.text = " 📱NOKIA 📱"
subtitle2.text = "reasons of failure"
img_path = "nokia_logo.png"
slide2.shapes.add_picture(img_path, Inches(2.5), Inches(5), Inches(5), Inches(1.5))

#introduction
slide_layout = prs.slide_layouts[5]
slide3 = prs.slides.add_slide(slide_layout)
title3 = slide3.shapes.title
fill3 = slide3.background.fill
fill3.solid()
fill3.fore_color.rgb = RGBColor(240, 248, 255)
title3.text = "Nokia’s fall is one of the most studied cases in business strategy and tech history."  \
"At its peak (mid-2000s), Nokia controlled more than 40% of the global mobile phone market."  \
"Within less than a decade, it lost dominance almost entirely." \
" Here’s a breakdown of why:"
p = title3.text_frame.paragraphs[0]
title3.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE   
run = p.runs[0]
run.font.size = Pt(25)        
run.font.bold = False         
run.font.name = "Arial"
title3.left = Inches(1)     # move left/right
title3.top = Inches(2)      # move up/down
title3.width = Inches(8)    # width of text box
title3.height = Inches(3)   # height of text box

#slide4
slide_layout = prs.slide_layouts[5]
slide4 = prs.slides.add_slide(slide_layout)
title4 = slide4.shapes.title
fill4 = slide4.background.fill
fill4.solid()
fill4.fore_color.rgb = RGBColor(204, 255, 204)
title4.text = (
    "1. Overconfidence & Complacency \n"
    "🔵Nokia was the global leader and thought its market share was untouchable. \n"
    "🔵Leadership underestimated the speed at which smartphones \n"
    "would disrupt the market.\n"
    "🔵Instead of adapting, Nokia doubled down on feature phones, assuming brand loyalty would protect them. \n"
"                                                                                                                   "
    "2. Software Weakness (Symbian vs. iOS/Android) \n"
    "🔵Nokia’s Symbian OS was powerful in its time but became clunky and outdated. \n"
    "🔵Developers found it difficult to build apps for Symbian, while iOS and Android created thriving app ecosystems. \n"
    "🔵Nokia failed to pivot quickly to a modern, user-friendly smartphone OS. \n"
)
p = title4.text_frame.paragraphs[0]
title4.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
for paragraph in title4.text_frame.paragraphs:
    paragraph.alignment = PP_ALIGN.LEFT  
run = p.runs[0]
for paragraph in title4.text_frame.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(20)   # set font size
        run.font.bold = False
        run.font.name = "Arial"
title4.left = Inches(1)     # move left/right
title4.top = Inches(2)      # move up/down
title4.width = Inches(8)    # width of text box
title4.height = Inches(3)   # height of text box

#slide5 
slide_layout = prs.slide_layouts[5]
slide5 = prs.slides.add_slide(slide_layout)
title5 = slide5.shapes.title
fill5 = slide5.background.fill
fill5.solid()
fill5.fore_color.rgb = RGBColor(255, 250, 205)
title5.text = (
    "3. Slow Decision-Making & Internal Politics \n"
    "🔵Nokia was plagued by bureaucracy and infighting between divisions. \n"
    "🔵Employees often knew the company was falling behind but leadership ignored or delayed critical decisions. \n"
    "🔵Innovation suffered due to risk-averse culture. \n"
    "                                                                                                                   "
    "4. Missed Opportunities \n"
    "🔵Nokia actually had touch-screen prototypes before the iPhone but dismissed them as impractical. \n"
    "🔵It also had early smartphone concepts with internet access and apps but never pushed them aggressively. \n"
)
p = title5.text_frame.paragraphs[0]
title5.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
for paragraph in title5.text_frame.paragraphs:
    paragraph.alignment = PP_ALIGN.LEFT  
run = p.runs[0]
for paragraph in title5.text_frame.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(20)   # set font size
        run.font.bold = False
        run.font.name = "Arial"
title5.left = Inches(1)     # move left/right
title5.top = Inches(2)      # move up/down
title5.width = Inches(8)    # width of text box
title5.height = Inches(3)   # height of text box

#slide 6
slide_layout = prs.slide_layouts[5]
slide6 = prs.slides.add_slide(slide_layout)
title6 = slide6.shapes.title
fill6 = slide6.background.fill
fill6.solid()
fill6.fore_color.rgb = RGBColor(255, 228, 225)
title6.text = (
    "5. Strategic Missteps \n"
    "🔵In 2011, Nokia partnered with Microsoft to adopt Windows Phone. While bold, it tied Nokia to an OS that never gained significant traction. \n"
    "🔵By then, Android was already dominant, and iOS was entrenched at the high end. Windows Phone became a “third option” that never scaled. \n"
    "                                                                                                                   "
    "6. Brand Perception Shift \n"
    "🔵Consumers began to see Nokia as “yesterday’s phone”; reliable but boring. \n"
    "🔵Apple and Samsung positioned themselves as innovative, stylish, and aspirational. \n"
    "🔵Nokia couldn’t rebrand itself fast enough."
)
p = title6.text_frame.paragraphs[0]
title6.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
for paragraph in title6.text_frame.paragraphs:
    paragraph.alignment = PP_ALIGN.LEFT  
run = p.runs[0]
for paragraph in title6.text_frame.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(20)   # set font size
        run.font.bold = False
        run.font.name = "Arial"
title6.left = Inches(1)     # move left/right
title6.top = Inches(2)      # move up/down
title6.width = Inches(8)    # width of text box
title6.height = Inches(3)   # height of text box

#slide7 
slide_layout = prs.slide_layouts[5]
slide7 = prs.slides.add_slide(slide_layout)
title7 = slide7.shapes.title
fill7 = slide7.background.fill
fill7.solid()
fill7.fore_color.rgb = RGBColor(224, 255, 255)
shape = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0, 120, 215)  # blue banner
shape.text = "📌 LESSON: Adaptability beats legacy!"
title7.text = ("Nokia didn’t fail because it lacked technology. It failed because it lacked adaptability. "
             "Its hardware was world class, but it underestimated the importance of software ecosystems "
             "and user experience, and it couldn’t align its organization fast enough to respond to the smartphone revolution."
)
p = title7.text_frame.paragraphs[0]
title7.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
for paragraph in title7.text_frame.paragraphs:
    paragraph.alignment = PP_ALIGN.LEFT  
run = p.runs[0]
for paragraph in title7.text_frame.paragraphs:
    for run in paragraph.runs:
        run.font.size = Pt(20)   # set font size
        run.font.bold = False
        run.font.name = "Arial"
title7.left = Inches(1)     # move left/right
title7.top = Inches(2)      # move up/down
title7.width = Inches(8)    # width of text box
title7.height = Inches(3)   # height of text box

#name slide
slide_layout = prs.slide_layouts[0]
slide8 = prs.slides.add_slide(slide_layout)
subtitle = slide8.placeholders[1]
title = slide8.shapes.title
fill8 = slide8.background.fill
fill8.solid()
fill8.fore_color.rgb = RGBColor(173,216,230) #light blue
title.text = "Made by:"
subtitle.text = "Lavanya Gupta"

#saving the ppt
prs.save("powerpynt_submission.pptx")

